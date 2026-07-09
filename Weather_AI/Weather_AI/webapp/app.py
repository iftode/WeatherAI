import os
import re
import threading
from io import BytesIO
from decimal import Decimal

import pandas as pd
from docx import Document
from flask import Flask, jsonify, render_template, request, send_file, session
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches

from config import Config
from db import connect_db, get_db_config_from_env, get_schema, ping, run_select, schema_to_text


CITY_ALIASES = {
    "bucuresti": "Bucharest",
    "bucurești": "Bucharest",
    "bucureşti": "Bucharest",
    "atena": "Athens",
    "roma": "Rome",
    "londra": "London",
    "viena": "Vienna",
    "varsovia": "Warsaw",
    "varșovia": "Warsaw",
    "varşovia": "Warsaw",
    "praga": "Prague",
    "bruxelles": "Brussels",
    "copenhaga": "Copenhagen",
    "lisabona": "Lisbon",
    "berna": "Bern",
    "dublin": "Dublin",
    "madrid": "Madrid",
    "paris": "Paris",
    "berlin": "Berlin",
    "budapesta": "Budapest",
    "sofia": "Sofia",
    "belgrad": "Belgrade",
    "bratislava": "Bratislava",
    "amsterdam": "Amsterdam",
    "stockholm": "Stockholm",
    "oslo": "Oslo",
    "helsinki": "Helsinki",
    "reykjavik": "Reykjavik",
    "ankara": "Ankara",
    "moscova": "Moscow",
    "kiev": "Kyiv",
    "chișinău": "Chisinau",
    "chişinău": "Chisinau",
    "chisinau": "Chisinau",
}

COUNTRY_ALIASES = {
    "romania": "Romania",
    "românia": "Romania",
    "germania": "Germany",
    "franta": "France",
    "franța": "France",
    "franţa": "France",
    "spania": "Spain",
    "italia": "Italy",
    "grecia": "Greece",
    "anglia": "United Kingdom",
    "marea britanie": "United Kingdom",
    "ungaria": "Hungary",
    "cehia": "Czechia",
    "austria": "Austria",
    "belgia": "Belgium",
    "olanda": "Netherlands",
    "portugalia": "Portugal",
    "irlanda": "Ireland",
    "polonia": "Poland",
    "bulgaria": "Bulgaria",
    "serbia": "Serbia",
    "slovacia": "Slovakia",
    "croatia": "Croatia",
    "croația": "Croatia",
    "croaţia": "Croatia",
    "suedia": "Sweden",
    "norvegia": "Norway",
    "finlanda": "Finland",
    "islanda": "Iceland",
    "turcia": "Turkey",
    "rusia": "Russia",
    "ucraina": "Ukraine",
    "moldova": "Moldova",
    "elveția": "Switzerland",
    "elveţia": "Switzerland",
    "elvetia": "Switzerland",
}

CONTINENT_ALIASES = {
    "europa": "Europe",
    "africa": "Africa",
    "asia": "Asia",
    "america de nord": "North America",
    "america de sud": "South America",
    "oceania": "Oceania",
    "australia": "Oceania",
}

def normalize_user_question(text: str) -> str:
    result = text

    for ro, en in sorted(CITY_ALIASES.items(), key=lambda x: len(x[0]), reverse=True):
        result = re.sub(rf"\b{re.escape(ro)}\b", en, result, flags=re.IGNORECASE)

    for ro, en in sorted(COUNTRY_ALIASES.items(), key=lambda x: len(x[0]), reverse=True):
        result = re.sub(rf"\b{re.escape(ro)}\b", en, result, flags=re.IGNORECASE)

    for ro, en in sorted(CONTINENT_ALIASES.items(), key=lambda x: len(x[0]), reverse=True):
        result = re.sub(rf"\b{re.escape(ro)}\b", en, result, flags=re.IGNORECASE)

    return result

app = Flask(__name__)
app.secret_key = Config.SECRET_KEY

OPENAI_MODEL = Config.OPENAI_MODEL
MAX_ROWS = Config.MAX_ROWS

api_key = Config.OPENAI_API_KEY
if not api_key:
    raise RuntimeError("OPENAI_API_KEY lipsește din fișierul .env")

client = OpenAI(api_key=api_key)

STATE = {
    "db_conn": None,
    "db_schema": None,
    "db_name": None,
    "schema_text": "",
    "conversation": [],
    "last_result": [],
    "last_sql": "",
    "language": "Română",
}


LANGUAGE_PROMPTS = {
    "Română": "Răspunde exclusiv în limba română.",
    "English": "Respond exclusively in English.",
    "Deutsch": "Antworte ausschließlich auf Deutsch.",
    "Français": "Réponds exclusivement en français.",
    "Español": "Responde exclusivamente en español.",
}


def build_system_prompt(schema_text: str, language: str) -> str:
    language_instruction = LANGUAGE_PROMPTS.get(language, LANGUAGE_PROMPTS["Română"])

    return f"""
Ești un asistent AI pentru interogarea unei baze de date meteorologice MySQL.

{language_instruction}

Reguli stricte:
1. Generezi DOAR interogări SELECT.
2. Nu folosești INSERT, UPDATE, DELETE, DROP, ALTER, TRUNCATE, CREATE.
3. Răspunsul tău trebuie să fie DOAR SQL, fără explicații.
4. Tabelul principal este weather_data.
5. Pentru orașe folosești coloana city.
6. Pentru țări folosești coloana country.
7. Pentru continente folosești coloana continent.
8. Pentru entitatea din fișierul sursă poți folosi coloana entity.
9. NU folosi niciodată condiții de forma entity = 'capital'.
10. Dacă utilizatorul cere Europa, Africa, Asia, America de Nord, America de Sud, Oceania, folosește coloana continent.
11. Dacă utilizatorul cere o comparație pe luni, folosește MONTH(date_record).
12. Dacă utilizatorul cere o comparație pe an, folosește YEAR(date_record).
13. Pentru topuri și comparații folosește AVG, MIN, MAX, SUM, GROUP BY, ORDER BY.
14. Dacă folosești funcții agregate și există filtre pe rezultatul agregat, folosește HAVING, nu WHERE.
15. Nu pune funcții agregate direct în WHERE.
16. Pentru denumiri românești folosește echivalentele internaționale:
   Bucuresti/București/Bucureşti = Bucharest
   Atena = Athens
   Roma = Rome
   Londra = London
   Viena = Vienna
   Varșovia/Varşovia = Warsaw
17. Generează întotdeauna SQL complet valid MySQL, cu SELECT ... FROM weather_data ...
18. Nu adăuga markdown, comentarii sau explicații.
19. Dacă nu există LIMIT, acesta va fi adăugat automat de aplicație.

Schema bazei de date:
{schema_text}
""".strip()


def build_sql_fix_prompt(schema_text: str, failed_sql: str, db_error: str, language: str) -> str:
    language_instruction = LANGUAGE_PROMPTS.get(language, LANGUAGE_PROMPTS["Română"])

    return f"""
Ești un asistent care repară SQL MySQL invalid.

{language_instruction}

Trebuie să returnezi DOAR SQL valid, fără explicații, fără markdown.

Reguli:
1. DOAR SELECT.
2. Tabelul principal este weather_data.
3. Dacă există agregări filtrate, folosește HAVING, nu WHERE.
4. NU folosi entity = 'capital'.
5. Pentru continente folosește coloana continent.
6. Pentru Bucuresti/București/Bucureşti folosește Bucharest.
7. Pentru Atena folosește Athens.
8. Pentru Roma folosește Rome.
9. Pentru Londra folosește London.
10. Păstrează intenția query-ului original.
11. Răspunsul trebuie să fie o singură interogare SQL validă MySQL.

Schema:
{schema_text}

SQL invalid:
{failed_sql}

Eroare MySQL:
{db_error}
""".strip()


def ensure_limit(sql: str, max_rows: int) -> str:
    cleaned = sql.strip()
    cleaned = cleaned.replace("```sql", "").replace("```", "").strip()
    cleaned = cleaned.rstrip(";").strip()

    aggregate_patterns = ["avg(", "min(", "max(", "sum(", "count("]
    has_aggregate = any(p in cleaned.lower() for p in aggregate_patterns)
    has_group_by = " group by " in cleaned.lower()

    if re.search(r"\blimit\s+\d+(\s*,\s*\d+)?\s*$", cleaned, re.IGNORECASE):
        return cleaned + ";"

    if has_aggregate and not has_group_by:
        return cleaned + ";"

    return f"{cleaned} LIMIT {max_rows};"


def validate_sql(sql: str):
    if not sql or not sql.strip():
        return False, "SQL gol."

    lowered = sql.lower().strip()
    forbidden = [
        "insert ", "update ", "delete ", "drop ", "alter ",
        "truncate ", "create ", "replace ", "grant ", "revoke "
    ]

    for word in forbidden:
        if word in lowered:
            return False, f"SQL interzis: {word.strip()}"

    if not lowered.startswith("select"):
        return False, "Este permis doar SELECT."

    if lowered.count(";") > 1:
        return False, "Este permisă o singură interogare."

    if "entity = 'capital'" in lowered or 'entity="capital"' in lowered:
        return False, "Query invalid: nu există valoarea generică entity='capital'."

    return True, ""


def auto_connect():
    config = get_db_config_from_env()
    conn = connect_db(config)
    schema = get_schema(conn, config.database)

    STATE["db_conn"] = conn
    STATE["db_schema"] = schema
    STATE["db_name"] = config.database
    STATE["schema_text"] = schema_to_text(schema)


def is_number(value):
    return isinstance(value, (int, float, Decimal)) and not isinstance(value, bool)


def to_float_safe(value):
    if value is None:
        return None
    if isinstance(value, Decimal):
        return float(value)
    if isinstance(value, (int, float)):
        return float(value)
    return None


def normalize_rows(rows):
    normalized = []
    for row in rows:
        normalized_row = {}
        for key, value in row.items():
            if isinstance(value, Decimal):
                normalized_row[key] = float(value)
            else:
                normalized_row[key] = value
        normalized.append(normalized_row)
    return normalized


def build_statistics(rows):
    if not rows:
        return {
            "rows_count": 0,
            "numeric_columns": [],
            "stats": {}
        }

    rows = normalize_rows(rows)
    columns = list(rows[0].keys())
    numeric_columns = []

    for col in columns:
        values = [to_float_safe(r.get(col)) for r in rows]
        values = [v for v in values if v is not None]
        if values:
            numeric_columns.append(col)

    stats = {}
    for col in numeric_columns:
        values = [to_float_safe(r.get(col)) for r in rows]
        values = [v for v in values if v is not None]

        if values:
            stats[col] = {
                "count": len(values),
                "min": round(min(values), 2),
                "max": round(max(values), 2),
                "avg": round(sum(values) / len(values), 2),
                "sum": round(sum(values), 2),
            }

    return {
        "rows_count": len(rows),
        "numeric_columns": numeric_columns,
        "stats": stats
    }


def get_result_dataframe():
    rows = STATE.get("last_result", [])
    if not rows:
        return pd.DataFrame()
    return pd.DataFrame(normalize_rows(rows))


def make_download_error(message: str, status_code: int = 400):
    return jsonify({"success": False, "message": message}), status_code


def execute_sql_with_repair(sql: str, language: str):
    try:
        rows = run_select(STATE["db_conn"], sql)
        return normalize_rows(rows), sql
    except Exception as first_error:
        fix_prompt = build_sql_fix_prompt(
            STATE["schema_text"],
            sql,
            str(first_error),
            language
        )

        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[{"role": "system", "content": fix_prompt}],
            temperature=0
        )

        fixed_sql = response.choices[0].message.content.strip()
        fixed_sql = ensure_limit(fixed_sql, MAX_ROWS)

        valid, error_message = validate_sql(fixed_sql)
        if not valid:
            raise RuntimeError(f"SQL reparat invalid: {error_message}. SQL: {fixed_sql}")

        rows = run_select(STATE["db_conn"], fixed_sql)
        return normalize_rows(rows), fixed_sql

def classify_question_intent(user_question: str, language: str) -> str:
    prompt = f"""
Clasifică întrebarea utilizatorului pentru o aplicație meteo.

Categorii posibile:
- weather_query: întrebarea cere date meteo, temperatură, umiditate, vânt, precipitații, orașe, țări, continente sau statistici meteo.
- conversation_history: utilizatorul cere istoricul conversației.
- unrelated: întrebarea nu are legătură cu aplicația meteo.

Returnează DOAR una dintre valorile:
weather_query
conversation_history
unrelated

Întrebare:
{user_question}
""".strip()

    response = client.chat.completions.create(
        model=OPENAI_MODEL,
        messages=[
            {"role": "system", "content": "Ești un clasificator. Răspunzi doar cu una dintre categoriile cerute."},
            {"role": "user", "content": prompt}
        ],
        temperature=0
    )

    intent = response.choices[0].message.content.strip().lower()

    if intent not in ["weather_query", "conversation_history", "unrelated"]:
        return "unrelated"

    return intent


def summarize_results_with_llm(user_question: str, rows: list, language: str) -> str:
    if not rows:
        return "Nu au fost găsite rezultate pentru această întrebare."

    df = pd.DataFrame(normalize_rows(rows))
    csv_data = df.head(50).to_csv(index=False)

    prompt = f"""
Primești rezultatele unei interogări SQL în format CSV.

Scrie o singură frază clară și naturală care rezumă rezultatele.
Nu inventa informații.
Dacă sunt puține date, descrie exact ce se observă.

Limba răspunsului: {language}

Întrebarea utilizatorului:
{user_question}

Rezultate CSV:
{csv_data}
""".strip()

    response = client.chat.completions.create(
        model=OPENAI_MODEL,
        messages=[
            {"role": "system", "content": "Ești un asistent care sumarizează rezultate meteo într-o singură frază."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.2
    )

    return response.choices[0].message.content.strip()


def save_chat_history(user_question, generated_sql=None, summary=None, answer_type="weather_query"):
    with STATE["db_conn"].cursor() as cursor:
        cursor.execute(
            """
            INSERT INTO chat_history
            (user_question, generated_sql, summary, answer_type)
            VALUES (%s, %s, %s, %s)
            """,
            (user_question, generated_sql, summary, answer_type)
        )
        STATE["db_conn"].commit()


def get_chat_history(limit=50):
    with STATE["db_conn"].cursor() as cursor:
        cursor.execute(
            """
            SELECT user_question, generated_sql, summary, answer_type, created_at
            FROM chat_history
            ORDER BY created_at DESC
            LIMIT %s
            """,
            (limit,)
        )
        return normalize_rows(cursor.fetchall())

@app.before_request
def init_app_once():
    if STATE["db_conn"] is None:
        auto_connect()


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/api/status", methods=["GET"])
def api_status():
    ok = STATE["db_conn"] is not None and ping(STATE["db_conn"])
    return jsonify({
        "connected": ok,
        "database": STATE["db_name"],
        "language": STATE["language"],
    })


@app.route("/api/set_language", methods=["POST"])
def api_set_language():
    data = request.get_json(force=True)
    language = (data.get("language") or "Română").strip()

    if language not in LANGUAGE_PROMPTS:
        language = "Română"

    STATE["language"] = language
    session["language"] = language

    return jsonify({
        "success": True,
        "language": language
    })


@app.route("/api/clear_history", methods=["POST"])
def api_clear_history():
    STATE["conversation"] = []
    STATE["last_result"] = []
    STATE["last_sql"] = ""
    session["conversation"] = []
    session.modified = True

    return jsonify({"success": True})


@app.route("/api/chat", methods=["POST"])
def api_chat():
    try:
        data = request.get_json(force=True)
        user_question = (data.get("message") or "").strip()

        if not user_question:
            return jsonify({"success": False, "message": "Mesajul este gol."}), 400

        if STATE["db_conn"] is None or not ping(STATE["db_conn"]):
            auto_connect()

        language = session.get("language", STATE["language"])
        STATE["language"] = language

        intent = classify_question_intent(user_question, language)

        if intent == "unrelated":
            answer_text = "Întrebarea nu are legătură cu aplicația meteo. Te rog să introduci o întrebare despre vreme, temperatură, precipitații, vânt, umiditate, orașe sau țări."

            save_chat_history(
                user_question=user_question,
                generated_sql=None,
                summary=answer_text,
                answer_type="unrelated"
            )

            return jsonify({
                "success": True,
                "question": user_question,
                "normalized_question": user_question,
                "sql": None,
                "rows": [],
                "columns": [],
                "message": answer_text,
                "summary": answer_text,
                "statistics": build_statistics([]),
                "language": language,
                "answer_type": "unrelated"
            })

        if intent == "conversation_history":
            history_rows = get_chat_history(limit=50)
            answer_text = "Am afișat istoricul conversației."

            save_chat_history(
                user_question=user_question,
                generated_sql=None,
                summary=answer_text,
                answer_type="conversation_history"
            )

            return jsonify({
                "success": True,
                "question": user_question,
                "normalized_question": user_question,
                "sql": None,
                "rows": history_rows,
                "columns": list(history_rows[0].keys()) if history_rows else [],
                "message": answer_text,
                "summary": answer_text,
                "statistics": build_statistics(history_rows),
                "language": language,
                "answer_type": "conversation_history"
            })

        normalized_question = normalize_user_question(user_question)

        conversation = session.get("conversation", [])
        conversation.append({"role": "user", "content": normalized_question})
        conversation = conversation[-6:]

        system_prompt = build_system_prompt(STATE["schema_text"], language)
        messages = [{"role": "system", "content": system_prompt}] + conversation

        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=messages,
            temperature=0
        )

        sql = response.choices[0].message.content.strip()
        sql = ensure_limit(sql, MAX_ROWS)

        valid, error_message = validate_sql(sql)
        if not valid:
            save_chat_history(
                user_question=user_question,
                generated_sql=sql,
                summary=error_message,
                answer_type="invalid_sql"
            )

            return jsonify({
                "success": False,
                "message": error_message,
                "sql": sql
            }), 400

        rows, final_sql = execute_sql_with_repair(sql, language)

        visible_rows = [r for r in rows if any(v is not None for v in r.values())]

        STATE["last_result"] = visible_rows
        STATE["last_sql"] = final_sql

        statistics = build_statistics(visible_rows)

        summary_text = summarize_results_with_llm(user_question, visible_rows, language)

        conversation.append({"role": "assistant", "content": summary_text})
        session["conversation"] = conversation
        session["language"] = language
        STATE["conversation"] = conversation

        save_chat_history(
            user_question=user_question,
            generated_sql=final_sql,
            summary=summary_text,
            answer_type="weather_query"
        )

        return jsonify({
            "success": True,
            "question": user_question,
            "normalized_question": normalized_question,
            "sql": final_sql,
            "rows": visible_rows,
            "columns": list(visible_rows[0].keys()) if visible_rows else [],
            "message": summary_text,
            "summary": summary_text,
            "statistics": statistics,
            "language": language,
            "answer_type": "weather_query"
        })

    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/api/chart_data", methods=["GET"])
def api_chart_data():
    rows = STATE.get("last_result", [])

    if not rows:
        return jsonify({"success": False, "message": "Nu există date pentru grafic."}), 400

    rows = [r for r in rows if any(v is not None for v in r.values())]
    if not rows:
        return jsonify({"success": False, "message": "Nu există date pentru grafic."}), 400

    columns = list(rows[0].keys())

    label_col = None
    value_col = None

    text_columns = []
    numeric_columns = []

    for col in columns:
        values = [r.get(col) for r in rows]

        if any(isinstance(v, str) and v.strip() for v in values if v is not None):
            text_columns.append(col)

        numeric_values = [to_float_safe(v) for v in values]
        numeric_values = [v for v in numeric_values if v is not None]
        if numeric_values:
            numeric_columns.append(col)

    if numeric_columns:
        value_col = numeric_columns[0]

    if text_columns:
        label_col = text_columns[0]

    if not label_col and value_col and len(rows) == 1:
        labels = [value_col]
        values = [to_float_safe(rows[0].get(value_col))]
        values = [0 if v is None else v for v in values]

        return jsonify({
            "success": True,
            "label_column": "metric",
            "value_column": value_col,
            "labels": labels,
            "values": values,
        })

    if not label_col or not value_col:
        return jsonify({"success": False, "message": "Nu s-au putut detecta coloanele pentru grafic."}), 400

    labels = [str(r.get(label_col, "")) for r in rows]
    values = [to_float_safe(r.get(value_col)) for r in rows]
    values = [0 if v is None else v for v in values]

    return jsonify({
        "success": True,
        "label_column": label_col,
        "value_column": value_col,
        "labels": labels,
        "values": values,
    })


@app.route("/api/statistics", methods=["GET"])
def api_statistics():
    rows = STATE.get("last_result", [])
    rows = [r for r in rows if any(v is not None for v in r.values())]

    if not rows:
        return jsonify({"success": False, "message": "Nu există date pentru statistici."}), 400

    return jsonify({
        "success": True,
        "statistics": build_statistics(rows)
    })


@app.route("/export/csv", methods=["GET"])
def export_csv():
    df = get_result_dataframe()

    if df.empty:
        return make_download_error("Nu există date pentru export.")

    output = BytesIO()
    csv_data = df.to_csv(index=False, encoding="utf-8-sig")
    output.write(csv_data.encode("utf-8-sig"))
    output.seek(0)

    return send_file(
        output,
        mimetype="text/csv",
        as_attachment=True,
        download_name="weather_results.csv"
    )


@app.route("/export/excel", methods=["GET"])
def export_excel():
    df = get_result_dataframe()

    if df.empty:
        return make_download_error("Nu există date pentru export.")

    output = BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        df.to_excel(writer, sheet_name="Results", index=False)

    output.seek(0)

    return send_file(
        output,
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        as_attachment=True,
        download_name="weather_results.xlsx"
    )


@app.route("/export/word", methods=["GET"])
def export_word():
    df = get_result_dataframe()

    if df.empty:
        return make_download_error("Nu există date pentru export.")

    doc = Document()
    doc.add_heading("Weather AI Agent - Results", level=1)
    doc.add_paragraph(f"Database: {STATE['db_name']}")
    doc.add_paragraph(f"SQL: {STATE['last_sql']}")

    table = doc.add_table(rows=1, cols=len(df.columns))
    table.style = "Table Grid"

    hdr_cells = table.rows[0].cells
    for i, col in enumerate(df.columns):
        hdr_cells[i].text = str(col)

    for _, row in df.iterrows():
        row_cells = table.add_row().cells
        for i, value in enumerate(row):
            row_cells[i].text = str(value)

    output = BytesIO()
    doc.save(output)
    output.seek(0)

    return send_file(
        output,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        as_attachment=True,
        download_name="weather_results.docx"
    )


@app.route("/export/pptx", methods=["GET"])
def export_pptx():
    df = get_result_dataframe()

    if df.empty:
        return make_download_error("Nu există date pentru export.")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    if slide.shapes.title:
        slide.shapes.title.text = "Weather AI Agent - Results"

    rows_count = min(len(df) + 1, 16)
    cols_count = len(df.columns)

    left = Inches(0.4)
    top = Inches(1.2)
    width = Inches(12.0)
    height = Inches(5.5)

    table_shape = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
    table = table_shape.table

    for i, col in enumerate(df.columns):
        table.cell(0, i).text = str(col)

    limited_df = df.head(rows_count - 1)
    for r_idx, (_, row) in enumerate(limited_df.iterrows(), start=1):
        for c_idx, value in enumerate(row):
            table.cell(r_idx, c_idx).text = str(value)

    output = BytesIO()
    prs.save(output)
    output.seek(0)

    return send_file(
        output,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name="weather_results.pptx"
    )


def open_browser():
    try:
        os.startfile("http://127.0.0.1:5000")
    except Exception as e:
        print("Nu am putut deschide browserul automat:", e)


if __name__ == "__main__":
    threading.Timer(1.5, open_browser).start()
    app.run(debug=True, use_reloader=False)