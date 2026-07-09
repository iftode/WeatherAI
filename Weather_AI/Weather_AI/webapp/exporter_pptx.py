from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt


def _add_textbox(slide, text, left, top, width, height, font_size=16):
    tx = slide.shapes.add_textbox(left, top, width, height).text_frame
    tx.word_wrap = True
    p = tx.paragraphs[0]
    p.text = text or ""
    p.font.size = Pt(font_size)
    return tx


def build_pptx(conversation, database: str = "") -> BytesIO:
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI Database Agent"
    slide.placeholders[1].text = (
        f"Conversație completă\n"
        f"DB: {database or ''}\n"
        f"Export: {datetime.now().isoformat(timespec='seconds')}"
    )

    for idx, item in enumerate(conversation, start=1):
        classification = item.get("classification", {}) or {}
        kind = classification.get("kind", "")
        intent = classification.get("intent", "")

        item_type = item.get("type", "") or ""
        question = item.get("question", "") or ""
        answer = item.get("answer", "") or ""
        explanation = item.get("explanation", "") or ""
        sql = item.get("sql")
        rows = item.get("rows") or []
        columns = item.get("columns") or []

        # Slide 1 for each interaction - summary
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Interacțiunea {idx}"

        summary = (
            f"Întrebare:\n{question}\n\n"
            f"Tip: {item_type}\n"
            f"Kind: {kind}\n"
            f"Intent: {intent}"
        )

        _add_textbox(
            slide,
            summary,
            Inches(0.6), Inches(1.0), Inches(12.0), Inches(1.9),
            font_size=16
        )

        # Text principal de afișat
        content_parts = []

        if sql:
            content_parts.append(f"SQL:\n{str(sql)[:1200]}")

        if explanation:
            content_parts.append(f"Explicație:\n{str(explanation)[:1200]}")
        elif answer:
            content_parts.append(f"Răspuns:\n{str(answer)[:1200]}")

        if not content_parts:
            content_parts.append("Nu există SQL sau explicație disponibilă pentru această interacțiune.")

        main_text = "\n\n".join(content_parts)

        _add_textbox(
            slide,
            main_text,
            Inches(0.6), Inches(3.0), Inches(12.0), Inches(3.2),
            font_size=16
        )

        # Slide 2 doar dacă există rezultate tabelare
        if columns and rows:
            slide2 = prs.slides.add_slide(prs.slide_layouts[5])
            slide2.shapes.title.text = f"Rezultate {idx} (preview)"

            preview_rows = rows[:10]
            rows_count = len(preview_rows) + 1
            cols_count = len(columns)

            left = Inches(0.4)
            top = Inches(1.3)
            width = Inches(12.5)
            height = Inches(4.8)

            table = slide2.shapes.add_table(rows_count, cols_count, left, top, width, height).table

            for c, col in enumerate(columns):
                table.cell(0, c).text = "" if col is None else str(col)

            for r_i, r in enumerate(preview_rows, start=1):
                for c_i, val in enumerate(r):
                    table.cell(r_i, c_i).text = "" if val is None else str(val)

            _add_textbox(
                slide2,
                f"Rânduri returnate total: {len(rows)} | În prezentare afișez primele {len(preview_rows)}",
                Inches(0.5), Inches(6.4), Inches(12.0), Inches(0.5),
                font_size=12
            )

        elif columns and not rows:
            slide2 = prs.slides.add_slide(prs.slide_layouts[5])
            slide2.shapes.title.text = f"Rezultate {idx} (preview)"

            _add_textbox(
                slide2,
                "Interogarea a returnat coloane, dar nu există rânduri de afișat.",
                Inches(0.6), Inches(1.5), Inches(12.0), Inches(1.0),
                font_size=18
            )

        elif item_type == "data":
            slide2 = prs.slides.add_slide(prs.slide_layouts[5])
            slide2.shapes.title.text = f"Rezultate {idx} (preview)"

            _add_textbox(
                slide2,
                "(Fără rezultate)",
                Inches(0.6), Inches(1.5), Inches(8.0), Inches(1.0),
                font_size=20
            )

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio