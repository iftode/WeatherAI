from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


SLIDE_W = 13.333
SLIDE_H = 7.5

COLOR_BLUE = RGBColor(47, 85, 151)
COLOR_BLUE_LIGHT = RGBColor(221, 235, 247)
COLOR_BLUE_MID = RGBColor(91, 155, 213)
COLOR_TEXT = RGBColor(40, 40, 40)
COLOR_TEXT_SOFT = RGBColor(90, 90, 90)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ROW_1 = RGBColor(242, 246, 252)
COLOR_ROW_2 = RGBColor(230, 238, 248)
COLOR_BORDER = RGBColor(200, 210, 225)


def normalize_value(val):
    if isinstance(val, list):
        return ", ".join(str(x) for x in val)
    if isinstance(val, tuple):
        return ", ".join(str(x) for x in val)
    if isinstance(val, set):
        return ", ".join(str(x) for x in val)
    if isinstance(val, dict):
        return str(val)
    if val is None:
        return ""
    return str(val)


def shorten_sql(sql_text, max_lines=5, max_len=420):
    sql_text = normalize_value(sql_text).strip()
    if not sql_text:
        return "-"
    lines = [line.rstrip() for line in sql_text.splitlines() if line.strip()]
    if len(lines) > max_lines:
        lines = lines[:max_lines] + ["..."]
    text = "\n".join(lines)
    if len(text) > max_len:
        text = text[: max_len - 3] + "..."
    return text


def shorten_text(text, max_len=320):
    text = normalize_value(text).strip()
    if not text:
        return "-"
    if len(text) <= max_len:
        return text
    return text[: max_len - 3] + "..."


def add_background(slide):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(250, 250, 250)
    bg.line.fill.background()

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W), Inches(0.75)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLOR_BLUE
    header.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.75), Inches(SLIDE_W), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_BLUE_MID
    accent.line.fill.background()


def add_title(slide, text, subtitle=None):
    add_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.12), Inches(10.8), Inches(0.35))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.68), Inches(0.43), Inches(8.0), Inches(0.18))
        tf = sub_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(10.5)
        p.font.color.rgb = RGBColor(225, 235, 248)


def add_card(slide, left, top, width, height):
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = COLOR_BORDER
    card.line.width = Pt(1)
    return card


def add_cover_slide(prs, database):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(11.6), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "AI Database Agent"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.15), Inches(11.2), Inches(1.4))
    tf = sub_box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = "Conversație completă exportată"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXT_SOFT

    p = tf.add_paragraph()
    p.text = f"Bază de date: {database or '-'}"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_SOFT

    p = tf.add_paragraph()
    p.text = f"Export: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_SOFT

    add_card(slide, Inches(1.1), Inches(1.0), Inches(11.1), Inches(3.0))


def add_label_value_line(tf, label, value, size=15, after=4):
    if len(tf.paragraphs) == 1 and not tf.paragraphs[0].text:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()

    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(after)

    r1 = p.add_run()
    r1.text = f"{label}: "
    r1.font.bold = True
    r1.font.size = Pt(size)
    r1.font.color.rgb = COLOR_TEXT

    r2 = p.add_run()
    r2.text = value if value else "-"
    r2.font.size = Pt(size)
    r2.font.color.rgb = COLOR_TEXT_SOFT


def add_text_slide(slide, title, question, tip, kind, intent, sql_text, explanation):
    add_title(slide, title, "Detalii interacțiune")

    add_card(slide, Inches(0.7), Inches(1.05), Inches(12.0), Inches(5.95))

    # Stânga: metadate
    meta_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.3), Inches(4.0), Inches(2.2))
    tf = meta_box.text_frame
    tf.clear()
    tf.word_wrap = True

    add_label_value_line(tf, "Întrebare", shorten_text(question, 120), 15, 8)
    add_label_value_line(tf, "Tip", tip, 14, 3)
    add_label_value_line(tf, "Kind", kind, 14, 3)
    add_label_value_line(tf, "Intent", intent, 14, 6)

    # Dreapta sus: SQL
    sql_title = slide.shapes.add_textbox(Inches(5.15), Inches(1.28), Inches(1.2), Inches(0.2))
    tf = sql_title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "SQL"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT

    sql_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.55), Inches(6.0), Inches(2.15)
    )
    sql_card.fill.solid()
    sql_card.fill.fore_color.rgb = COLOR_BLUE_LIGHT
    sql_card.line.color.rgb = RGBColor(180, 200, 225)

    sql_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.72), Inches(5.65), Inches(1.8))
    tf = sql_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = shorten_sql(sql_text, max_lines=5)
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_TEXT

    # Jos: explicație
    exp_title = slide.shapes.add_textbox(Inches(0.95), Inches(4.0), Inches(2.0), Inches(0.2))
    tf = exp_title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Explicație"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT

    exp_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(4.28), Inches(10.9), Inches(1.8)
    )
    exp_card.fill.solid()
    exp_card.fill.fore_color.rgb = RGBColor(248, 248, 248)
    exp_card.line.color.rgb = COLOR_BORDER

    exp_box = slide.shapes.add_textbox(Inches(1.15), Inches(4.46), Inches(10.5), Inches(1.4))
    tf = exp_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = shorten_text(explanation, 420)
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_SOFT


def choose_columns(columns):
    excluded = {
        "description",
        "special_features",
        "full_text",
        "payload",
        "raw",
        "query",
        "sql",
    }
    filtered = [c for c in columns if c not in excluded]
    return filtered[:6]


def get_column_widths(display_columns):
    n = len(display_columns)

    if n == 1:
        return [11.55]
    if n == 2:
        return [6.8, 4.75]
    if n == 3:
        return [2.3, 3.2, 6.05]
    if n == 4:
        return [1.7, 2.4, 1.8, 5.65]
    if n == 5:
        return [1.5, 1.7, 1.5, 1.7, 5.15]
    if n == 6:
        if "payment_date" in display_columns:
            return [1.3, 1.3, 1.2, 1.3, 1.3, 5.15]
        return [1.3, 2.2, 1.2, 1.4, 2.0, 3.45]

    return [11.55 / n] * n


def add_table_slide(slide, title, columns, rows, total_rows=None):
    add_title(slide, title, "Preview rezultate")
    add_card(slide, Inches(0.72), Inches(1.02), Inches(12.0), Inches(5.98))

    display_columns = choose_columns(columns)
    col_index_map = [columns.index(c) for c in display_columns]

    display_rows = []
    for row in rows[:10]:
        if isinstance(row, dict):
            display_rows.append([normalize_value(row.get(c, "")) for c in display_columns])
        else:
            display_rows.append(
                [normalize_value(row[i]) if i < len(row) else "" for i in col_index_map]
            )

    n_rows = len(display_rows) + 1
    n_cols = len(display_columns)

    table_left = Inches(0.95)
    table_top = Inches(1.42)
    table_width = Inches(11.55)
    table_height = Inches(4.50)

    shape = slide.shapes.add_table(
        n_rows, n_cols, table_left, table_top, table_width, table_height
    )
    table = shape.table

    widths = get_column_widths(display_columns)
    for i, w in enumerate(widths[:n_cols]):
        table.columns[i].width = Inches(w)

    # Header
    for c, col_name in enumerate(display_columns):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_BLUE
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT

        r = p.add_run()
        r.text = col_name
        r.font.bold = True
        r.font.size = Pt(11.5)
        r.font.color.rgb = COLOR_WHITE

    # Rows
    for r_idx, row in enumerate(display_rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_ROW_1 if r_idx % 2 else COLOR_ROW_2
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]

            try:
                float(str(value).replace(",", "."))
                is_number = True
            except Exception:
                is_number = False

            p.alignment = PP_ALIGN.CENTER if is_number else PP_ALIGN.LEFT

            text = str(value)
            if display_columns[c_idx] != "payment_date" and len(text) > 34:
                text = text[:31] + "..."

            r = p.add_run()
            r.text = text
            r.font.size = Pt(10.5)
            r.font.color.rgb = COLOR_TEXT

    note_box = slide.shapes.add_textbox(Inches(1.0), Inches(6.08), Inches(10.7), Inches(0.25))
    tf = note_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    if total_rows is None:
        p.text = f"În prezentare sunt afișate primele {len(display_rows)} rânduri."
    else:
        p.text = f"Rânduri returnate total: {total_rows} | În prezentare sunt afișate primele {len(display_rows)}."
    p.font.size = Pt(10.5)
    p.font.color.rgb = COLOR_TEXT_SOFT
    p.alignment = PP_ALIGN.LEFT


def build_pptx(conversation, database: str = "") -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    add_cover_slide(prs, database)

    for idx, item in enumerate(conversation, start=1):
        classification = item.get("classification", {}) or {}
        question = normalize_value(item.get("question", ""))
        tip = normalize_value(classification.get("type", classification.get("tip", "")))
        kind = normalize_value(classification.get("kind", ""))
        intent = normalize_value(classification.get("intent", ""))
        sql_text = normalize_value(item.get("sql", ""))
        explanation = normalize_value(item.get("answer", ""))
        result = item.get("result")

        # slide text
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_text_slide(
            slide,
            f"Interacțiunea {idx}",
            question,
            tip,
            kind,
            intent,
            sql_text,
            explanation,
        )

        # slide rezultate
        if isinstance(result, dict) and result.get("columns") and result.get("rows") is not None:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_table_slide(
                slide,
                f"Rezultate {idx}",
                result.get("columns", []),
                result.get("rows", []),
                total_rows=result.get("total_rows", len(result.get("rows", []))),
            )

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio